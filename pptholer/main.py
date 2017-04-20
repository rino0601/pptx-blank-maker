# coding=utf-8
import io

import click
from pptx import Presentation

__author__ = 'rino0601'


@click.command()
@click.argument('source_pptx_file_path', type=click.Path(exists=True))
def cli(source_pptx_file_path):
    click.echo(u"Welcome, Before we start converting, There is some things you have to know.")
    click.echo(u"   1. This tool does not support videos which you upload tistory directly.\n"
               u"   They seems only play able on tistory blog.\n"
               u"   If you have good idea for this, please contribute this project.\n"
               u"   2. This tool ignore comments, track backs, tags and categories. \n"
               u"   There is some different with jekyll and tistory's system.\n"
               u"   I didn't have enough time for consider that.\n"
               u"   3. If you want bring your picture and attachment from tistory, that blog has to be alive.\n"
               u"   (backup_with_attachment.xml isn't available for now.)\n")
    # click.confirm(u'Do you want to continue?', abort=True)
    # with click.progressbar(posts, label=u"Converting posts...") as bar:
    file_path = click.format_filename(source_pptx_file_path)
    with io.open(file_path, 'rb') as source_pptx_file:
        prs = Presentation(source_pptx_file)
        with io.open(file_path + '.txt', 'w', encoding='utf8') as notepad:
            for s_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for i, paragraph in enumerate(shape.text_frame.paragraphs):
                        text_row = "".join([run.text for run in paragraph.runs])
                        print(s_num, i, text_row)
                        notepad.write(text_row)
                        notepad.write('\n')

                        # for run in paragraph.runs:
                        #     run.text += "ㅗ"

    click.echo("DONE")
