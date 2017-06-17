# coding=utf-8
import glob
import io
import os

import click
from pptx import Presentation

__author__ = 'rino0601'


@click.command()
@click.argument('source_pptx_file_path', type=click.Path(exists=True))
def cli(source_pptx_file_path):
    origin_dir = os.getcwd()
    click.echo(u"Start converting...")
    # click.confirm(u'Do you want to continue?', abort=True)
    # with click.progressbar(posts, label=u"Converting posts...") as bar:

    filename = click.format_filename(source_pptx_file_path)
    file_path_list = [filename]
    if os.path.isdir(filename):
        click.echo(u"Searching *.pptx files in given directory...")
        os.chdir(source_pptx_file_path)
        file_path_list = glob.glob('*.pptx')
        click.echo(file_path_list)

    for file_path in file_path_list:
        click.echo(u"Processing %r ..." % file_path)
        with io.open(file_path, 'rb') as source_pptx_file:
            prs = Presentation(source_pptx_file)
            with io.open(file_path + '.txt', 'w', encoding='utf8') as notepad:
                for s_num, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for i, paragraph in enumerate(shape.text_frame.paragraphs):
                            text_row = "".join([run.text for run in paragraph.runs])
                            # click.echo(s_num, i, text_row)
                            notepad.write(text_row)
                            notepad.write('\n')

                            # for run in paragraph.runs:
                            #     run.text += "ㅗ"

    os.chdir(origin_dir)
    click.echo("DONE")

if __name__ == '__main__':
    cli()
